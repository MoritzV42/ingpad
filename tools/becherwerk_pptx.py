#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Erzeugt die PowerPoint-Ausarbeitung zur Lernsituation 3 (Becherwerk).

Rendert die gesamte Loesung (Rahmendaten + Teilaufgaben A-I + Antriebsstrang-
Uebersicht) aus dem ingpad-Canvas in eine projektier-/druckbare .pptx.

    python becherwerk_pptx.py [ZIEL.pptx]

Designprinzip: helles, sachliches Theme mit ingpad-Akzent (Blau/Violett),
eine Folie je Teilaufgabe nach dem Schema Gegeben / Gesucht / Ansatz / Ergebnis.
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Farben (ingpad-Brand) -------------------------------------------------
BG = RGBColor(0xF4, 0xF6, 0xF9)
INK = RGBColor(0x1B, 0x25, 0x30)
MUT = RGBColor(0x5A, 0x6B, 0x7B)
ACC = RGBColor(0x1B, 0x6F, 0xC4)   # ingpad-Blau (kraeftig fuer hellen Grund)
VIO = RGBColor(0x7A, 0x5A, 0xF0)   # Violett-Akzent
GREEN = RGBColor(0x1F, 0x9D, 0x57)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD2, 0xDA, 0xE3)
DARK = RGBColor(0x0F, 0x16, 0x1F)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tb, tf


def para(tf, text, size=14, color=INK, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, space_after=4, bullet=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    pre = "•  " if bullet else ""
    r = p.add_run(); r.text = pre + text
    f = r.font; f.size = Pt(size); f.name = FONT
    f.bold = bold; f.italic = italic; f.color.rgb = color
    return p


def rect(s, x, y, w, h, fill, line=None, line_w=0.75):
    sp = s.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def header(s, letter, title, accent=ACC):
    """Kopfbalken mit Buchstaben-Badge + Titel."""
    bar = rect(s, 0, 0, EMU_W, Inches(1.05), accent)
    # Badge
    bd = rect(s, Inches(0.45), Inches(0.2), Inches(0.65), Inches(0.65), CARD)
    tf = bd.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = letter
    r.font.size = Pt(26); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = accent
    # Titel
    _, tf2 = box(s, Inches(1.3), Inches(0.18), Inches(11.6), Inches(0.7))
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf2, title, size=22, color=CARD, bold=True, first=True)
    return bar


def panel(s, x, y, w, h, label, lines, label_color=ACC, result=False):
    """Beschriftetes Feld (Gegeben/Gesucht/Ansatz/Ergebnis)."""
    fill = CARD if not result else RGBColor(0xEA, 0xF3, 0xFD)
    ln = LINE if not result else ACC
    rect(s, x, y, w, h, fill, line=ln, line_w=1.0 if result else 0.75)
    _, tf = box(s, x + Pt(8), y + Pt(6), w - Pt(16), h - Pt(12))
    para(tf, label, size=11, color=label_color if not result else ACC, bold=True,
         first=True, space_after=5)
    for it in lines:
        if isinstance(it, tuple):
            txt, kw = it
        else:
            txt, kw = it, {}
        para(tf, txt, size=kw.get("size", 13), color=kw.get("color", INK),
             bold=kw.get("bold", False), italic=kw.get("italic", False),
             bullet=kw.get("bullet", False), space_after=kw.get("sa", 3))


def footer(s, n):
    _, tf = box(s, Inches(0.45), Inches(7.05), Inches(8), Inches(0.35))
    para(tf, "Modul 3 · Technische Lösungen entwickeln · LS 3 — Antriebsstrang Becherwerk",
         size=9, color=MUT, first=True)
    _, tf2 = box(s, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
    para(tf2, str(n), size=9, color=MUT, first=True, align=PP_ALIGN.RIGHT)


# ============================ TITELFOLIE ====================================
s = slide(DARK)
rect(s, 0, Inches(3.05), EMU_W, Inches(0.06), ACC)
_, tf = box(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2))
para(tf, "Antriebsstrang für ein Becherwerk", size=44, color=CARD, bold=True, first=True)
_, tf = box(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(0.9))
para(tf, "Entwicklung und Auslegung des kompletten Antriebsstrangs", size=20,
     color=RGBColor(0x8A, 0xA0, 0xB4), first=True)
_, tf = box(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.6))
para(tf, "Modul 3 · Technische Lösungen entwickeln", size=16, color=ACC, bold=True, first=True, space_after=6)
para(tf, "Lernsituation 3 — Vollständige Handlung A–I", size=15, color=RGBColor(0xC0, 0xCE, 0xDC))
para(tf, "Becher · Antriebsleistung · Motor · Übersetzungen · Stirnrad · Kettentrieb · "
         "Gleitlager · Wälzlager · Kupplung", size=12, color=RGBColor(0x8A, 0xA0, 0xB4))
_, tf = box(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5))
para(tf, "erstellt mit ingpad — github.com/MoritzV42/ingpad", size=11,
     color=RGBColor(0x6A, 0x7B, 0x8B), first=True)

# ========================== RAHMENDATEN =====================================
s = slide()
header(s, "0", "Rahmendaten & Aufgabenstellung", accent=VIO)
col_w = Inches(6.05)
panel(s, Inches(0.45), Inches(1.3), col_w, Inches(2.7), "FÖRDERUNG — VORGABEN", [
    ("Förderleistung  V̇ = 60 m³/h", {}),
    ("Dichte (max.)  ρ = 1800 kg/m³", {}),
    ("max. Geschwindigkeit  v_max = 2 m/s", {}),
    ("Becherabstand  a = 500 mm", {}),
    ("Förderhöhe  h ≈ 15 m", {}),
    ("Betriebsdauer  12 h/Tag", {}),
    ("Leistungszuschlag  +10 %", {}),
], label_color=VIO)
panel(s, Inches(6.83), Inches(1.3), col_w, Inches(2.7), "RANDBEDINGUNGEN ANTRIEBSSTRANG", [
    ("2 Stirnradstufen (i ≤ 6/Stufe, x = 0)", {}),
    ("Kettentrieb als Hülltrieb  i_K = 1…2", {}),
    ("Zahnkranz-Ø  D ≈ 630 mm", {}),
    ("Becherkette  m′ = 10,25 kg/m", {}),
    ("Wellen-Ø (Getriebe) 60 mm / Hauptwelle 70 mm", {}),
    ("Zähnezahlreihe: 16,18,20,25,28,35,45,50,56,63,71,80,90,100", {"size": 11, "color": MUT}),
])
panel(s, Inches(0.45), Inches(4.2), Inches(12.43), Inches(2.5), "ERGEBNISSE IM ÜBERBLICK", [
    ("Becher 4,75 L · v = 1,75 m/s · P_erf ≈ 4,86 kW · η_ges ≈ 0,85", {"bold": True, "color": ACC}),
    ("Motor SEW DRN132M4 · 7,5 kW · n = 1468 1/min · M_N = 49 Nm", {}),
    ("Übersetzung gesamt i_ges = 27,6  (Soll 27,7):  Stufe 1 = 20/90 · Stufe 2 = 18/71 · Kette = 18/28", {}),
    ("Stirnrad St.1: m = 4 mm (ungeh.) / m = 2,5 mm (geh.) · Kette 16B-1, a ≈ 570 mm", {}),
    ("Gleitlager Rotguss (p·v ≈ 0,8) · Wälzlager 6014-2RSR (L10h ≈ 8,6·10⁵ h) · Kupplung ROTEX, M_erf ≈ 86 Nm", {}),
])
footer(s, 2)

# ============================ AUFGABEN A-I ==================================
TASKS = [
    dict(letter="A", title="Elevatorbecher auswählen",
         given=["V̇ = 60 m³/h", "v_max = 2 m/s", "Becherabstand a = 500 mm"],
         sought=["Bechergröße (Volumen)", "Fördergeschwindigkeit v", "grafische Kontrolle"],
         approach=["Volumenstrom:  V̇ = (v / a) · V_Becher", "→ Becher aus Reihe, v rückrechnen"],
         result=[("Becher 4,75 L (DIN 15234, tiefe Form)", {"bold": True}),
                 ("v = 60 / 4,75 · 7,2 = 1,75 m/s", {}),
                 ("Maße: b = 315 · a = 180 · h₁ = 224 · h₂ = 118 · r = 56 mm", {"size": 11, "color": MUT}),
                 ("im Volumenstrom-Diagramm (a = 500 mm) bestätigt: 60 m³/h ↔ 4,75 L ↔ 1,75 m/s", {"size": 11, "color": MUT})]),
    dict(letter="B", title="Antriebsleistung ermitteln",
         given=["V̇ = 60 m³/h", "ρ = 1800 kg/m³", "h ≈ 15 m", "Zuschlag +10 %"],
         sought=["erforderliche Leistung P_erf", "an der Becherwerkswelle"],
         approach=["Massenstrom (Becher/Kette heben sich auf):  ṁ = V̇ · ρ = 30 kg/s",
                   "Hubleistung:  P_Hub = ṁ · g · h", "P_erf = P_Hub · 1,1"],
         result=[("P_Hub = 30 · 9,81 · 15 = 4414,5 W ≈ 4,41 kW", {}),
                 ("P_erf = 4,41 · 1,1 = 4,86 kW", {"bold": True})]),
    dict(letter="C", title="Antriebsmotor auswählen",
         given=["P_erf = 4,86 kW", "η_ges = η_Verz · η_Kette · η_Wälzl · η_Gleitl ≈ 0,85"],
         sought=["erforderliche Motorleistung", "Motorauswahl (Katalog SEW DRN)"],
         approach=["P_Mot ≥ P_erf / η_ges"],
         result=[("P_Mot = 4,86 / 0,85 = 5,72 kW", {}),
                 ("→ 5,5 kW reicht nicht → nächste Baugröße 7,5 kW", {"color": MUT}),
                 ("SEW DRN132M4:  7,5 kW · n_N = 1468 1/min · M_N = 49 Nm · η = 90,4 %", {"bold": True})]),
    dict(letter="D", title="Übersetzungen & Zähnezahlen",
         given=["n_Mot = 1468 1/min", "v = 1,75 m/s · Zahnkranz-Ø D = 630 mm",
                "i_K = 1…2 · 2 Stirnradstufen i ≤ 6 · x = 0"],
         sought=["n_HW · i_ges", "Aufteilung i_K, i₁, i₂", "Zähnezahlen z₁…z₆ + n/M je Welle"],
         approach=["n_HW = v / (π·D) · 60", "i_ges = n_Mot / n_HW = i₁ · i₂ · i_K", "i = z₂ / z₁"],
         result=[("n_HW = 1,75 / (π·0,630) · 60 = 53,05 1/min  →  i_ges = 1468 / 53,05 = 27,67", {}),
                 ("Stufe 1: 20/90 (i₁ = 4,50) · Stufe 2: 18/71 (i₂ = 3,94) · Kette: 18/28 (i_K = 1,56)", {"bold": True}),
                 ("Kontrolle: i_ist = 4,50 · 3,94 · 1,56 = 27,61 (−0,22 %) → n_HW = 53,2 1/min", {"size": 11, "color": MUT}),
                 ("Wellen:  W1 1468/5,27kW/34,3Nm · W2 326/5,17/151 · W3 83/5,06/585 · HW 53/4,86/873", {"size": 11, "color": MUT})]),
    dict(letter="E", title="Stirnrad Stufe 1: Modul + Festigkeitsnachweis",
         given=["z₁ = 20 · z₂ = 90 · u = 4,5", "M₁ = 34,3 Nm · K_A = 1,4 · gerade · x = 0",
                "ψ = 0,8 (ungeh.) / 0,6 (geh.)"],
         sought=["erf. Modul m", "Nachweis Zahnfuß σ_Fb ≤ σ_Flim", "Nachweis Zahnflanke σ_Hp ≤ σ_Hlim"],
         approach=["Auslegungsmoment  M_t = K_A · M₁ = 1,4 · 34,3 = 48,0 Nm",
                   "Modul aus Flanken-/Fußtragfähigkeit, anschließend Nachweis"],
         result=[("Ungehärtet C45E N:  m = 4 mm · b = 64 mm  (S_F ≈ 9,5 · S_H ≈ 2,3 → Flanke maßgebend)", {"bold": True}),
                 ("Gehärtet 15CrNi6:  m = 2,5 mm · b = 30 mm  (S_F ≈ 3,9 · S_H ≈ 3,2)", {"bold": True}),
                 ("→ gehärtete Variante deutlich kompakter (d₁ halbiert) — Empfehlung für Serie", {"size": 11, "color": MUT})]),
    dict(letter="F", title="Kettentrieb (Hauptkette zum Becherwerk)",
         given=["n₃ = 82,7 1/min · z₅ = 18 · z₆ = 28 (i_K = 1,556)",
                "P₃ = 5,06 kW · K_A = 1,4 · a₀ ≈ 560 mm"],
         sought=["korrigierte Leistung P_c", "Kettengröße (Teilung)", "Glieder X + Achsabstand a"],
         approach=["P_c = P · K_A · f_s ,  f_s = (19/z₅)^1,08",
                   "X = 2a₀/p + (z₅+z₆)/2 + ((z₆−z₅)/2π)² · p/a₀"],
         result=[("f_s = (19/18)^1,08 = 1,060  →  P_c = 5,06 · 1,4 · 1,060 = 7,51 kW", {}),
                 ("Kette 16B-1 (p = 25,4 mm, 1″) · X = 68 Glieder · a ≈ 570 mm (+1,8 %)", {"bold": True}),
                 ("z₅ = 18 ≥ 17 → ruhiger Lauf, geringer Polygoneffekt", {"size": 11, "color": MUT})]),
    dict(letter="G", title="Gleitlager Zwischenwelle",
         given=["d = 60 mm · n₂ = 326,2 1/min · M₂ = 151,2 Nm",
                "Ritzel St.2: m ≈ 4 · z₃ = 18 → d₃ = 72 mm · α = 20° · 2 Lager"],
         sought=["Lagerkraft", "Lagerbreite b (b/d)", "Werkstoff: p und p·v prüfen", "Ölversorgung"],
         approach=["F_t = 2M₂/d₃ · F_r = F_t·tanα · F_res = √(F_t² + F_r²)",
                   "p = F/(d·b) ≤ p_zul · v = d·π·n / 60000"],
         result=[("F_t = 4200 N · F_r = 1529 N · F_res = 4470 N → je Lager 2235 N", {}),
                 ("b = 48 mm (b/d = 0,8) · Rotguss G-CuSn10Zn", {"bold": True}),
                 ("p = 0,78 N/mm² · p·v = 0,80 ≪ Grenzwerte → große Reserve, hydrodynamisch", {}),
                 ("Schmiernut in Lastfreizone · Tauch-/Spritzschmierung · Getriebeöl CLP 220", {"size": 11, "color": MUT})]),
    dict(letter="H", title="Wälzlager 6014-2RSR Hauptwelle",
         given=["6014-2RSR: C_r = 40500 N · C_0r = 31000 N · d = 70 mm",
                "n_HW = 53,2 1/min · P_erf = 4,86 kW · v = 1,75 m/s",
                "Becherkette m′ = 10,25 kg/m · h = 15 m"],
         sought=["Lagerlast aus Kettenzug + Gewicht", "L10h (dynamisch, p = 3)", "Bewertung"],
         approach=["F_u = P_erf / v", "L10 = (C/P)³ [Mio. U]", "L10h = L10·10⁶ / (60·n)"],
         result=[("F_u = 2777 N · G = 3017 N → F ≈ 5794 N → je Lager P ≈ 2897 N", {}),
                 ("L10 = (40500/2897)³ = 2733 Mio. U → L10h ≈ 8,6·10⁵ h (≈ 856 000 h)", {"bold": True}),
                 ("üblich für Stetigförderer 20 000–50 000 h → Lager weit darüber (steifigkeitsbestimmt)", {"size": 11, "color": MUT})]),
    dict(letter="I", title="Kupplung Motor ↔ Getriebe",
         given=["SEW DRN132M4: M_N = 49 Nm · n = 1468 1/min",
                "Last Becherwerk (Stetigförderer) · S_A ≈ 1,5…1,75"],
         sought=["Kupplungstyp", "Auslegungsmoment M_erf", "Einbau-/Montagehinweis"],
         approach=["M_erf = M_N · S_A ≤ M_KN  (Kupplungs-Nennmoment)"],
         result=[("M_erf = 49 · 1,75 = 85,8 Nm  (konservativ; bei S_A = 1,5: 73,5 Nm)", {}),
                 ("drehelastische Klauenkupplung mit Elastomer-Zahnkranz (KTR ROTEX GS 28, M_KN ≥ 90 Nm)", {"bold": True}),
                 ("dämpft Drehschwingungen · gleicht Wellenversatz aus · axiale Steckmontage, wartungsarm", {"size": 11, "color": MUT}),
                 ("Montage: Naben aufschieben, Spaltmaß E lt. Datenblatt, Ausrichtung mit Messuhr prüfen", {"size": 11, "color": MUT})]),
]

for i, t in enumerate(TASKS):
    s = slide()
    header(s, t["letter"], t["title"])
    cw = Inches(6.05)
    panel(s, Inches(0.45), Inches(1.25), cw, Inches(1.95), "GEGEBEN",
          [(g, {"bullet": True, "size": 13}) for g in t["given"]])
    panel(s, Inches(6.83), Inches(1.25), cw, Inches(1.95), "GESUCHT",
          [(g, {"bullet": True, "size": 13}) for g in t["sought"]], label_color=VIO)
    panel(s, Inches(0.45), Inches(3.35), Inches(12.43), Inches(1.4), "ANSATZ",
          [(a, {"size": 13, "italic": True}) for a in t["approach"]])
    panel(s, Inches(0.45), Inches(4.95), Inches(12.43), Inches(1.95),
          "ERGEBNIS " + t["letter"], t["result"], result=True)
    footer(s, i + 3)

# ===================== ANTRIEBSSTRANG-ÜBERSICHT =============================
s = slide(DARK)
_, tf = box(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.8))
para(tf, "Antriebsstrang — Zusammenfassung", size=30, color=CARD, bold=True, first=True)
rect(s, Inches(0.6), Inches(1.35), Inches(12.1), Pt(2), ACC)

# Kette als Text
_, tf = box(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.9))
para(tf, "Motor 1468 1/min  →  Kupplung ROTEX  →  Stirnrad 20/90  →  Zwischenwelle  →  "
         "Stirnrad 18/71  →  Kettenrad 18  →  Kette 16B-1  →  Hauptwelle 28 (53 1/min)",
     size=13, color=RGBColor(0xC0, 0xCE, 0xDC), first=True)

# Wellen-Tabelle
rows = [
    ["Welle", "n [1/min]", "P [kW]", "M [Nm]", "Element"],
    ["W1 · Eingang (Motor/Ritzel St.1)", "1468", "5,27", "34,3", "z₁ = 20"],
    ["W2 · Zwischenwelle", "326,2", "5,17", "151,2", "z₂ = 90 / z₃ = 18 · Gleitlager"],
    ["W3 · Ausgang / kl. Kettenrad", "82,7", "5,06", "584,5", "z₄ = 71 / z₅ = 18"],
    ["HW · Becherwerks-Hauptwelle", "53,2", "4,86", "872,9", "z₆ = 28 · Wälzlager 6014"],
]
tx, ty = Inches(0.6), Inches(2.7)
tw = Inches(12.1)
colw = [Inches(4.0), Inches(1.7), Inches(1.3), Inches(1.4), Inches(3.7)]
rh = Inches(0.62)
for ri, row in enumerate(rows):
    cx = tx
    head = ri == 0
    for ci, cell in enumerate(row):
        cellfill = ACC if head else (RGBColor(0x16, 0x21, 0x2C) if ri % 2 else RGBColor(0x1E, 0x29, 0x35))
        rc = rect(s, cx, ty + Emu(int(rh) * ri), colw[ci], rh, cellfill,
                  line=RGBColor(0x2C, 0x3A, 0x4A), line_w=0.5)
        tf = rc.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(6)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 or ci == 4 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = cell
        r.font.size = Pt(12 if not head else 12.5); r.font.name = FONT
        r.font.bold = head; r.font.color.rgb = CARD
        cx = cx + colw[ci]

_, tf = box(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7))
para(tf, "Motor-Betriebsmoment 34,3 Nm < M_N = 49 Nm → Teillastbetrieb mit Reserve · "
         "Gesamtwirkungsgrad η_ges ≈ 0,85 · i_ges = 27,6 (Abweichung zum Soll −0,2 %)",
     size=12, color=GREEN, first=True)

# ---- speichern -------------------------------------------------------------
ZIEL = sys.argv[1] if len(sys.argv) > 1 else "Becherwerk_Ausarbeitung.pptx"
prs.save(ZIEL)
print("gespeichert:", ZIEL, "·", len(prs.slides.__iter__.__self__._sldIdLst), "Folien")
